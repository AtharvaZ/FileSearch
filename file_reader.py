import os
import pdfplumber
from docx import Document
import hashlib
import pptx
import time

CODE_FILE_EXTENSIONS = {".py", ".js", ".java", ".cpp", ".c", ".ts", ".jsx", ".tsx", ".go", ".rs", ".rb"}
TEXT_FILE_EXTENSIONS = {".txt", ".md"}

def _get_file_metadata(file_path: str) -> tuple[str, str]:
    """Extract common file metadata (name and modified time)"""
    file_name = os.path.basename(file_path)
    modified_time = time.ctime(os.path.getmtime(file_path))
    return file_name, modified_time

def read_file(file_path: str) -> list[str]:
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == ".pdf":
        return return_pdf_info(file_path)
    elif file_ext == ".docx":
        return return_docx_info(file_path)
    elif file_ext in TEXT_FILE_EXTENSIONS or file_ext in CODE_FILE_EXTENSIONS:
        return return_text_info(file_path)
    elif file_ext == ".pptx":
        return return_pptx_info(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_ext}")

def return_pdf_info(file_path: str) -> list[str]:
    """Handle PDF files using pdfplumber (2-3x faster than PyPDF2)"""
    file_name, modified_time = _get_file_metadata(file_path)

    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                except Exception:
                    continue
    except Exception as e:
        raise Exception(f"Failed to read PDF: {str(e)}")

    if not text.strip():
        raise Exception("PDF contains no extractable text")

    hashed_content = hash_file(text)

    return [file_path, file_name, hashed_content, modified_time, text]

def return_docx_info(file_path: str) -> list[str]:
    """Handle DOCX files"""
    file_name, modified_time = _get_file_metadata(file_path)

    try:
        doc = Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    except Exception as e:
        raise Exception(f"Failed to read DOCX: {str(e)}")

    if not text.strip():
        raise Exception("DOCX contains no extractable text")

    hashed_content = hash_file(text)

    return [file_path, file_name, hashed_content, modified_time, text]

def return_text_info(file_path: str) -> list[str]:
    """Handle plain text files (.txt, .md, and code files)"""
    file_name, modified_time = _get_file_metadata(file_path)

    try:
        with open(file_path, "r", encoding="utf-8") as file:
            text = file.read()
    except Exception as e:
        raise Exception(f"Failed to read text file: {str(e)}")

    if not text.strip():
        raise Exception("File contains no extractable text")

    hashed_content = hash_file(text)

    return [file_path, file_name, hashed_content, modified_time, text]

def return_pptx_info(file_path: str) -> list[str]:
    file_name, modified_time = _get_file_metadata(file_path)

    try:
        presentation = pptx.Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        raise Exception(f"Failed to read PPTX: {str(e)}")

    if not text.strip():
        raise Exception("PPTX contains no extractable text")

    hashed_content = hash_file(text)

    return [file_path, file_name, hashed_content, modified_time, text]

def hash_file(text: str) -> str:
    return hashlib.sha256(text.encode()).hexdigest()
